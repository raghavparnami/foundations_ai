"""Port of src/lib/agent/generate-presentation.ts — the `generate_presentation` tool.

The TS implementation builds a native PPTX via `pptxgenjs`. We don't have
`python-pptx` in `backend/pyproject.toml` and the task forbids adding deps, so
the binary generation is stubbed: we still validate the spec, persist a
markdown narrative of the slides into `reports` (so the user can read it in
the UI), and return a 501-style error string for the .pptx itself. Swapping
in `python-pptx` later is a drop-in replacement for `_render_pptx_bytes`.
"""
from __future__ import annotations

import base64
import io
import json
import logging
import re
import time
from typing import Any

from psycopg.rows import dict_row
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

from ..audit import audit
from ..db import get_conn

log = logging.getLogger(__name__)


_SLIDE_TYPES = ("title", "summary", "chart", "table", "bullets", "closing")


TOOLS: list[dict[str, Any]] = [
    {
        "type": "function",
        "function": {
            "name": "generate_presentation",
            "description": (
                "Build a presentation-quality deck from a structured spec (title / summary / chart / "
                "table / bullets / closing slides). The binary .pptx exporter is a stub in this build "
                "— the markdown narrative IS persisted and rendered in the UI."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "spec": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "subtitle": {"type": "string"},
                            "author": {"type": "string"},
                            "slides": {
                                "type": "array",
                                "items": {"type": "object"},
                                "minItems": 1,
                            },
                        },
                        "required": ["title", "slides"],
                    },
                    "slug": {
                        "type": "string",
                        "description": "Optional stable slug. Defaults to a sluggified title.",
                    },
                },
                "required": ["spec"],
                "additionalProperties": False,
            },
        },
    },
]


async def run_tool(name: str, args: dict[str, Any], *, conversation_id: str) -> str:
    """Dispatch for generate_presentation. Returns a JSON string."""
    try:
        if name == "generate_presentation":
            return _json(
                await generate_presentation(
                    spec=args["spec"],
                    slug=args.get("slug"),
                    conversation_id=conversation_id,
                )
            )
        return _json({"error": f"Unknown tool: {name}"})
    except Exception as e:  # noqa: BLE001
        log.exception("generate_presentation.run_tool failed: %s", name)
        return _json({"error": f"{type(e).__name__}: {e}"})


async def generate_presentation(
    *,
    spec: dict[str, Any],
    slug: str | None,
    conversation_id: str,
) -> dict[str, Any]:
    """Validate the spec, persist the markdown narrative, return the slug.

    The actual .pptx generation is stubbed (see module docstring); the
    return value still mirrors the TS shape with `pptx_status='unavailable'`
    so the caller can detect the gap.
    """
    if not isinstance(spec, dict):
        return {"ok": False, "error": "spec is required"}
    title = (spec.get("title") or "").strip()
    if not title:
        return {"ok": False, "error": "spec.title is required"}
    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        return {"ok": False, "error": "spec.slides must be a non-empty array"}

    # Validate each slide has a known type. Per-type field requirements are
    # checked permissively — bad slides degrade to "(unsupported)" in the
    # narrative rather than refusing the whole deck.
    for i, s in enumerate(slides):
        if not isinstance(s, dict) or s.get("type") not in _SLIDE_TYPES:
            return {
                "ok": False,
                "error": (
                    f"spec.slides[{i}].type must be one of: {', '.join(_SLIDE_TYPES)} "
                    f"(got {s.get('type') if isinstance(s, dict) else type(s).__name__})"
                ),
            }

    final_slug = _sanitize_slug(slug or title) or f"deck-{int(time.time() * 1000)}"
    filename = f"{final_slug}.pptx"

    narrative = _render_markdown_narrative(spec)

    # Build the actual .pptx via python-pptx and store it as base64 in reports
    # under the .pptx slug. The /api/presentations/{slug}/download route
    # base64-decodes the body and streams the binary back.
    pptx_status = "ok"
    pptx_reason: str | None = None
    pptx_b64 = ""
    try:
        pptx_b64 = _build_pptx_b64(spec)
    except Exception as e:  # noqa: BLE001
        log.warning("generate_presentation.pptx_build_failed err=%s", e)
        pptx_status = "unavailable"
        pptx_reason = f"pptx build failed: {e}. Markdown narrative is still available."

    async with get_conn() as conn:
        async with conn.cursor(row_factory=dict_row) as cur:
            # Markdown narrative under the bare slug
            await cur.execute(
                """
                INSERT INTO reports (slug, title, body_md, conversation_id)
                     VALUES (%s, %s, %s, %s)
                ON CONFLICT (slug) DO UPDATE
                   SET title = EXCLUDED.title,
                       body_md = EXCLUDED.body_md,
                       conversation_id = EXCLUDED.conversation_id,
                       created_at = now()
                RETURNING id, slug
                """,
                (final_slug, title, narrative, conversation_id),
            )
            row = await cur.fetchone()
            assert row is not None

            # Binary .pptx (base64) under the .pptx slug, if built
            if pptx_status == "ok":
                await cur.execute(
                    """
                    INSERT INTO reports (slug, title, body_md, conversation_id)
                         VALUES (%s, %s, %s, %s)
                    ON CONFLICT (slug) DO UPDATE
                       SET title = EXCLUDED.title,
                           body_md = EXCLUDED.body_md,
                           conversation_id = EXCLUDED.conversation_id,
                           created_at = now()
                    """,
                    (filename, title, pptx_b64, conversation_id),
                )

    await audit(
        "agent",
        "generate_presentation",
        filename,
        {
            "conversationId": conversation_id,
            "slides": len(slides),
            "bytes": len(narrative),
            "pptx_status": pptx_status,
        },
    )
    return {
        "ok": True,
        "slug": filename,
        "title": title,
        "slide_count": len(slides),
        "download_url": f"/api/presentations/{final_slug}/download",
        "markdown_url": f"/api/reports/{filename}/download",
        "bytes": len(narrative),
        "pptx_status": pptx_status,
        "pptx_reason": pptx_reason,
    }


# ─── Markdown rendering (stand-in for the native PPTX) ──────────────────


def _render_markdown_narrative(spec: dict[str, Any]) -> str:
    title = (spec.get("title") or "").strip()
    subtitle = (spec.get("subtitle") or "").strip()
    author = (spec.get("author") or "").strip()
    slides = spec.get("slides") or []

    parts: list[str] = [f"# {title}"]
    if subtitle:
        parts += ["", f"_{subtitle}_"]
    if author:
        parts += ["", f"Prepared by {author}"]
    parts += [""]

    for i, slide in enumerate(slides, start=1):
        parts += _render_slide_md(i, slide)
        parts += [""]
    return "\n".join(parts).rstrip() + "\n"


def _render_slide_md(index: int, slide: dict[str, Any]) -> list[str]:
    t = slide.get("type")
    headline = (slide.get("headline") or slide.get("title") or "").strip()
    header = f"## Slide {index} · {t}: {headline}" if headline else f"## Slide {index} · {t}"
    out: list[str] = ["---", header, ""]

    if t == "title":
        sub = (slide.get("subtitle") or "").strip()
        if sub:
            out += [sub]
    elif t == "summary":
        body = (slide.get("body") or "").strip()
        if body:
            out += [body, ""]
        stats = slide.get("stats") or []
        for stat in stats:
            label = (stat.get("label") or "").strip()
            value = (stat.get("value") or "").strip()
            delta = (stat.get("delta") or "").strip()
            row = f"- **{label}**: {value}"
            if delta:
                row += f" ({delta})"
            out.append(row)
    elif t == "chart":
        out += [
            f"_Chart_: {slide.get('chart_type', 'bar')} "
            f"({slide.get('x_field')} vs {slide.get('y_field')})",
            "",
        ]
        data = slide.get("data") or []
        if isinstance(data, list) and data:
            keys = sorted({k for row in data for k in (row.keys() if isinstance(row, dict) else [])})
            if keys:
                out.append("| " + " | ".join(keys) + " |")
                out.append("|" + " --- |" * len(keys))
                for row in data[:50]:
                    if isinstance(row, dict):
                        out.append("| " + " | ".join(str(row.get(k, "")) for k in keys) + " |")
        cap = (slide.get("caption") or "").strip()
        if cap:
            out += ["", f"_{cap}_"]
    elif t == "table":
        cols = slide.get("columns") or []
        rows = slide.get("rows") or []
        if cols:
            out.append("| " + " | ".join(str(c) for c in cols) + " |")
            out.append("|" + " --- |" * len(cols))
        for row in rows[:100]:
            if isinstance(row, list):
                out.append("| " + " | ".join(str(c) for c in row) + " |")
        cap = (slide.get("caption") or "").strip()
        if cap:
            out += ["", f"_{cap}_"]
    elif t == "bullets":
        for b in slide.get("bullets") or []:
            out.append(f"- {b}")
    elif t == "closing":
        body = (slide.get("body") or "").strip()
        if body:
            out += [body]
    else:
        out += ["_(unsupported slide type)_"]
    return out


# ─── Native PPTX builder via python-pptx ─────────────────────────────────


def _build_pptx_b64(spec: dict[str, Any]) -> str:
    """CXO-level deck. Clean typography, generous whitespace, native charts.

    Returns a base64-encoded .pptx so it can be persisted as TEXT in the
    reports table (the /api/presentations/{slug}/download route decodes).
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title = (spec.get("title") or "Loom Deck").strip()
    subtitle = (spec.get("subtitle") or "").strip()
    author = (spec.get("author") or "").strip()

    _add_title_slide(prs, title, subtitle, author)
    for slide in spec.get("slides") or []:
        if not isinstance(slide, dict):
            continue
        st = slide.get("type")
        if st == "title":
            _add_title_slide(
                prs,
                (slide.get("headline") or slide.get("title") or "").strip(),
                (slide.get("subtitle") or "").strip(),
                author,
            )
        elif st == "summary":
            _add_summary_slide(prs, slide)
        elif st == "chart":
            _add_chart_slide(prs, slide)
        elif st == "table":
            _add_table_slide(prs, slide)
        elif st == "bullets":
            _add_bullets_slide(prs, slide)
        elif st == "closing":
            _add_closing_slide(prs, slide)

    buf = io.BytesIO()
    prs.save(buf)
    return base64.b64encode(buf.getvalue()).decode("ascii")


_NAVY = (20, 21, 42)
_INK = (35, 38, 65)
_MUTED = (91, 96, 117)
_ACCENT = (91, 108, 255)


def _add_title_slide(prs: "Presentation", title: str, subtitle: str, author: str) -> None:
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(12), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*_NAVY)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(*_MUTED)
    if author:
        sb = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(12), Inches(0.5))
        sp = sb.text_frame.paragraphs[0]
        sp.text = author
        sp.font.size = Pt(12)
        sp.font.color.rgb = RGBColor(*_MUTED)


def _add_section_header(slide: object, headline: str) -> None:
    from pptx.dml.color import RGBColor

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.8))  # type: ignore[attr-defined]
    p = tb.text_frame.paragraphs[0]
    p.text = headline
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*_NAVY)


def _add_summary_slide(prs: "Presentation", slide_spec: dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    headline = (slide_spec.get("headline") or "Summary").strip()
    _add_section_header(slide, headline)

    body = (slide_spec.get("body") or "").strip()
    if body:
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(1.5))
        p = tb.text_frame.paragraphs[0]
        p.text = body
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(*_INK)
        tb.text_frame.word_wrap = True

    stats = slide_spec.get("stats") or []
    if isinstance(stats, list) and stats:
        col_w = 11.5 / max(1, min(len(stats), 4))
        for i, stat in enumerate(stats[:4]):
            if not isinstance(stat, dict):
                continue
            left = Inches(0.7 + i * col_w)
            box = slide.shapes.add_textbox(left, Inches(3.2), Inches(col_w - 0.2), Inches(2))
            tf = box.text_frame
            tf.word_wrap = True
            p1 = tf.paragraphs[0]
            p1.text = str(stat.get("value") or "").strip() or "—"
            p1.font.size = Pt(40)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(*_ACCENT)
            p2 = tf.add_paragraph()
            p2.text = str(stat.get("label") or "").strip()
            p2.font.size = Pt(13)
            p2.font.color.rgb = RGBColor(*_MUTED)
            delta = str(stat.get("delta") or "").strip()
            if delta:
                p3 = tf.add_paragraph()
                p3.text = delta
                p3.font.size = Pt(11)
                p3.font.color.rgb = RGBColor(*_MUTED)


def _add_chart_slide(prs: "Presentation", slide_spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, (slide_spec.get("headline") or "Chart").strip())

    chart_type = (slide_spec.get("chart_type") or "bar").lower()
    x_field = str(slide_spec.get("x_field") or "x")
    y_field = str(slide_spec.get("y_field") or "y")
    data_rows = slide_spec.get("data") or []
    if not (isinstance(data_rows, list) and data_rows):
        return

    cats = [str(r.get(x_field, "")) for r in data_rows if isinstance(r, dict)]
    vals: list[float] = []
    for r in data_rows:
        if not isinstance(r, dict):
            continue
        try:
            vals.append(float(r.get(y_field, 0) or 0))
        except (TypeError, ValueError):
            vals.append(0.0)

    chart_data = CategoryChartData()
    chart_data.categories = cats
    chart_data.add_series(y_field, vals)

    chart_kind = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "area": XL_CHART_TYPE.AREA,
    }.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    slide.shapes.add_chart(
        chart_kind,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(5),
        chart_data,
    )


def _add_table_slide(prs: "Presentation", slide_spec: dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, (slide_spec.get("headline") or "Table").strip())

    cols = slide_spec.get("columns") or []
    rows = slide_spec.get("rows") or []
    if not cols:
        return
    n_rows = min(len(rows) + 1, 12)
    n_cols = min(len(cols), 8)
    tbl = slide.shapes.add_table(
        n_rows,
        n_cols,
        Inches(0.7),
        Inches(1.5),
        Inches(12),
        Inches(5),
    ).table

    for j in range(n_cols):
        cell = tbl.cell(0, j)
        cell.text = str(cols[j])
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(*_NAVY)

    for i in range(1, n_rows):
        row = rows[i - 1] if i - 1 < len(rows) else []
        if not isinstance(row, list):
            continue
        for j in range(n_cols):
            cell = tbl.cell(i, j)
            cell.text = str(row[j]) if j < len(row) else ""
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(*_INK)


def _add_bullets_slide(prs: "Presentation", slide_spec: dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_header(slide, (slide_spec.get("headline") or "Key points").strip())

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(slide_spec.get("bullets") or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(*_INK)
        p.space_after = Pt(12)


def _add_closing_slide(prs: "Presentation", slide_spec: dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    headline = (slide_spec.get("headline") or "Thank you").strip()
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(3), Inches(12), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = headline
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*_NAVY)
    body = (slide_spec.get("body") or "").strip()
    if body:
        tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(4), Inches(12), Inches(2))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = body
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(*_MUTED)
        tb2.text_frame.word_wrap = True


_SLUG_RE = re.compile(r"^[a-z0-9][a-z0-9-]{1,80}$")


def _sanitize_slug(raw: str) -> str | None:
    s = raw.strip().lower()
    s = re.sub(r"\s+", "-", s)
    s = re.sub(r"[^a-z0-9-]", "", s)
    s = re.sub(r"-+", "-", s)
    s = s.strip("-")
    return s if _SLUG_RE.match(s) else None


def _json(value: Any) -> str:
    return json.dumps(value, default=str)
